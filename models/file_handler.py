#model\file_handler.py
import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import pytesseract
from PIL import Image, ImageEnhance, UnidentifiedImageError
from io import BytesIO
import magic  # python-magic-bin
from typing import Optional, Union
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Constants
UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), '..', 'uploads')
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'png', 'jpg', 'jpeg'}

# Configure Tesseract path (Windows)
TESSERACT_PATH = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
if os.path.exists(TESSERACT_PATH):
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH

def allowed_file(filename: str) -> bool:
    """Check if the file has an allowed extension."""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(file: Union[BytesIO, str]) -> Optional[str]:
    """Extract text from supported files with dual detection (MIME + extension)."""
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    
    try:
        # Primary detection: MIME type
        file.seek(0)
        file_type = magic.from_buffer(file.read(1024), mime=True)
        file.seek(0)
        
        handlers = {
            'application/pdf': extract_text_from_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': extract_text_from_docx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': extract_text_from_pptx
        }
        
        if handler := handlers.get(file_type):
            return handler(file)
        elif file_type in ('image/png', 'image/jpeg'):
            return extract_text_from_image(file)
            
    except Exception as mime_error:
        logger.warning(f"MIME detection failed: {mime_error}")
        
    try:
        # Fallback: File extension
        if hasattr(file, 'filename'):
            filename = file.filename.lower()
            ext_handlers = {
                '.pdf': extract_text_from_pdf,
                '.docx': extract_text_from_docx,
                '.pptx': extract_text_from_pptx,
                '.png': extract_text_from_image,
                '.jpg': extract_text_from_image,
                '.jpeg': extract_text_from_image
            }
            for ext, handler in ext_handlers.items():
                if filename.endswith(ext):
                    return handler(file)
                    
    except Exception as ext_error:
        logger.error(f"Extension fallback failed: {ext_error}")
    
    return None

def extract_text_from_pdf(file: BytesIO) -> Optional[str]:
    """Extract text from PDF with PyMuPDF."""
    try:
        file.seek(0)
        with fitz.open(stream=file.read(), filetype="pdf") as doc:
            return "\n".join(page.get_text() for page in doc if page.get_text())
    except Exception as e:
        logger.error(f"PDF extraction failed: {e}")
        return None

def extract_text_from_docx(file: BytesIO) -> Optional[str]:
    """Extract text from DOCX."""
    try:
        file.seek(0)
        doc = Document(BytesIO(file.read()))
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception as e:
        logger.error(f"DOCX extraction failed: {e}")
        return None

def extract_text_from_pptx(file: BytesIO) -> Optional[str]:
    """Extract text from PPTX."""
    try:
        file.seek(0)
        prs = Presentation(BytesIO(file.read()))
        return "\n".join(
            shape.text.strip()
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
    except Exception as e:
        logger.error(f"PPTX extraction failed: {e}")
        return None

def extract_text_from_image(file: BytesIO) -> Optional[str]:
    """Enhanced OCR text extraction with image preprocessing."""
    try:
        file.seek(0)
        img = Image.open(BytesIO(file.read()))
        
        # Preprocessing for better OCR
        img = img.convert('L')  # Grayscale
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(2.0)  # Increase contrast
        
        # OCR with configured Tesseract
        text = pytesseract.image_to_string(img, config='--psm 6')
        return text.strip() if text else None
        
    except UnidentifiedImageError:
        logger.error("Invalid image file")
    except Exception as e:
        logger.error(f"OCR processing failed: {e}")
    
    return None